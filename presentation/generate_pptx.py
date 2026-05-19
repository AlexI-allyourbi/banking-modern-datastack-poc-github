"""Generate the training PowerPoint presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette ---
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT_BLUE = RGBColor(0x41, 0x9C, 0xF0)
ACCENT_GREEN = RGBColor(0x6B, 0xCB, 0x77)
ACCENT_YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
ACCENT_RED = RGBColor(0xF0, 0x65, 0x65)
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return tf

def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5), title, 44, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1), subtitle, 22, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    return slide

def section_slide(step_num, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    # Step badge
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(2), Inches(0.6), f"Step {step_num}", 14, ACCENT_YELLOW, True)
    # Title
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(12), Inches(0.8), title, 32, WHITE, True)
    # Divider
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    # Bullets
    add_bullet_slide(slide, Inches(0.7), Inches(1.9), Inches(11), Inches(4.5), bullets, 18, LIGHT_GRAY)
    if note:
        add_text_box(slide, Inches(0.7), Inches(6.5), Inches(11), Inches(0.6), note, 13, MED_GRAY)
    return slide

# ============================================================
# SLIDE 1: Title
# ============================================================
s = title_slide(
    "Banking Modern Data Stack",
    "End-to-End Data Engineering Training\nPostgreSQL  →  MinIO  →  Snowflake  →  dbt  →  Airflow"
)
add_text_box(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "AISMR  |  Training Workshop", 16, MED_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Agenda / Overview
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text_box(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), "Training Agenda", 36, ACCENT_BLUE, True)
items = [
    "Step 1-2:   Set up PostgreSQL & create database tables",
    "Step 3:       Generate fake banking data with Python Faker",
    "Step 4:       Set up Apache Airflow for orchestration",
    "Step 5:       Set up MinIO (S3-compatible object storage)",
    "Step 6-7:   Set up Snowflake & create raw tables",
    "Step 8-9:   Build & run the MinIO → Snowflake Airflow DAG",
    "Step 10-11: Initialize dbt & run transformations",
    "Step 12:     Run the SCD2 Snapshots DAG",
    "Step 13:     Verify data in Snowflake",
]
add_bullet_slide(s, Inches(0.7), Inches(1.5), Inches(11), Inches(5.5), items, 20, LIGHT_GRAY)

# ============================================================
# SLIDE 3: Architecture
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text_box(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), "Architecture Overview", 36, ACCENT_BLUE, True)
arch_lines = [
    "┌─────────────┐    ┌──────────┐    ┌───────────┐    ┌─────────────┐    ┌───────────┐",
    "│   Faker     │───→│ Postgres │───→│   MinIO   │───→│  Snowflake  │───→│    dbt    │",
    "│ (Data Gen)  │    │  (OLTP)  │    │  (S3)     │    │  (DW)       │    │(Transform)│",
    "└─────────────┘    └──────────┘    └───────────┘    └─────────────┘    └───────────┘",
    "",
    "                                Orchestrated by Apache Airflow",
    "",
    "Snowflake Layers:   RAW (VARIANT)  →  Staging (Views)  →  Snapshots (SCD2)  →  Marts (Dims/Facts)",
]
add_bullet_slide(s, Inches(0.5), Inches(1.5), Inches(12), Inches(5), arch_lines, 15, ACCENT_GREEN)
add_text_box(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5), "All infrastructure runs in Docker containers via docker-compose.yml", 14, MED_GRAY)

# ============================================================
# SLIDE 4: Tech Stack
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text_box(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), "Tech Stack", 36, ACCENT_BLUE, True)
items = [
    "PostgreSQL 15 — Source OLTP relational database (customers, accounts, transactions)",
    "Python + Faker — Synthetic data generation to simulate real banking activity",
    "MinIO — S3-compatible object storage (local data lake for Parquet files)",
    "Apache Kafka + Debezium — Real-time Change Data Capture (CDC) streaming",
    "Snowflake — Cloud data warehouse (RAW → Staging → Marts)",
    "dbt (data build tool) — SQL transformations, SCD2 snapshots, testing",
    "Apache Airflow — Workflow orchestration and scheduling (DAGs)",
    "Docker + docker-compose — Containerized infrastructure for all services",
]
add_bullet_slide(s, Inches(0.7), Inches(1.5), Inches(11), Inches(5.5), items, 18, LIGHT_GRAY)

# ============================================================
# SLIDE 5: Step 1-2 — PostgreSQL Setup
# ============================================================
section_slide("1-2", "Set Up PostgreSQL & Create Tables", [
    "PostgreSQL runs as a Docker container (image: postgres:15)",
    "WAL level set to 'logical' for Change Data Capture support",
    "Credentials managed via .env file (POSTGRES_USER, POSTGRES_PASSWORD, POSTGRES_DB)",
    "Data persisted to ./docker/postgres/data volume mount",
    "",
    "Three tables defined in postgres/schema.sql:",
    "    • customers — id, first_name, last_name, email (UNIQUE), created_at",
    "    • accounts — id, customer_id (FK), account_type, balance (CHECK ≥ 0), currency",
    "    • transactions — id, account_id (FK), txn_type, amount (CHECK > 0), related_account_id, status",
    "",
    "Performance index on transactions(account_id, created_at)",
], "Command: docker-compose up -d postgres  |  Then run schema.sql against the database")

# ============================================================
# SLIDE 6: Step 3 — Faker Data Generator
# ============================================================
section_slide("3", "Generate Fake Banking Data (Faker)", [
    "Script: data-generator/faker_generator.py",
    "Uses the Python Faker library to create realistic synthetic data",
    "",
    "Each iteration generates:",
    "    • 10 customers with random names & unique emails",
    "    • 2 accounts per customer (SAVINGS or CHECKING), random balance $10-$1000",
    "    • 50 transactions (DEPOSIT / WITHDRAWAL / TRANSFER), amounts up to $1000",
    "",
    "Runs in a continuous loop (2-second intervals) or single iteration (--once flag)",
    "Connects to PostgreSQL using psycopg2 with credentials from .env",
], "Command: python data-generator/faker_generator.py")

# ============================================================
# SLIDE 7: Step 4 — Airflow Setup
# ============================================================
section_slide("4", "Set Up Apache Airflow", [
    "Custom Docker image built from dockerfile-airflow.dockerfile:",
    "    • Base: apache/airflow:2.9.3",
    "    • Installs: dbt-core, dbt-snowflake",
    "",
    "Two containers: airflow-webserver (UI on port 8080) + airflow-scheduler",
    "",
    "Key volume mounts:",
    "    • ./docker/dags → /opt/airflow/dags (DAG files)",
    "    • ./banking_dbt → /opt/airflow/banking_dbt (dbt project)",
    "    • ./banking_dbt/.dbt → /home/airflow/.dbt (dbt profiles.yml)",
    "",
    "Separate airflow-postgres database on port 5433 for Airflow metadata",
], "Access UI: http://localhost:8080")

# ============================================================
# SLIDE 8: Step 5 — MinIO Setup
# ============================================================
section_slide("5", "Set Up MinIO (S3-Compatible Storage)", [
    "MinIO runs as a Docker container (image: minio/minio:latest)",
    "API on port 9000, Console UI on port 9001",
    "",
    "Acts as a local data lake — stores Parquet files between PostgreSQL and Snowflake",
    "",
    "Bucket structure:",
    "    raw/customers/customers_2026-05-15_143022.parquet",
    "    raw/accounts/accounts_2026-05-15_143022.parquet",
    "    raw/transactions/transactions_2026-05-15_143022.parquet",
    "",
    "Script: consumer/postgres_to_minio.py",
    "    → Reads all rows from PostgreSQL, converts to Parquet, uploads to MinIO",
], "Console: http://localhost:9001  |  Command: python consumer/postgres_to_minio.py")

# ============================================================
# SLIDE 9: Step 6-7 — Snowflake Setup
# ============================================================
section_slide("6-7", "Set Up Snowflake & Create Raw Tables", [
    "Sign up for Snowflake 30-day free trial at signup.snowflake.com",
    "",
    "Run these queries in Snowflake (from snowflake/queries.txt):",
    "    CREATE DATABASE banking;",
    "    CREATE SCHEMA banking.raw;",
    "    CREATE TABLE customers (v VARIANT);",
    "    CREATE TABLE accounts (v VARIANT);",
    "    CREATE TABLE transactions (v VARIANT);",
    "",
    "Why VARIANT? — Stores semi-structured data (Parquet/JSON) without rigid schema",
    "    → Fields extracted later in dbt staging models: v:id::string, v:email::string",
], "VARIANT = flexible schema for data lake ingestion pattern")

# ============================================================
# SLIDE 10: Step 8-9 — MinIO to Snowflake DAG
# ============================================================
section_slide("8-9", "MinIO → Snowflake Airflow DAG", [
    "DAG: minio_to_snowflake_banking (Schedule: every 1 minute)",
    "File: docker/dags/minio_to_snowflake_dag.py",
    "",
    "Task 1 — download_minio:",
    "    → Connects to MinIO via boto3 (S3 client)",
    "    → Lists & downloads all Parquet files to /tmp/minio_downloads",
    "",
    "Task 2 — load_snowflake:",
    "    → Receives file paths from Task 1 via XCom (cross-task communication)",
    "    → PUT files to Snowflake internal stage (@%tablename)",
    "    → COPY INTO table with FILE_FORMAT=(TYPE=PARQUET)",
    "",
    "Task flow: download_minio → load_snowflake",
], "Enable the DAG in Airflow UI → it auto-runs every minute")

# ============================================================
# SLIDE 11: Step 10-11 — dbt Setup & Run
# ============================================================
section_slide("10-11", "Initialize dbt & Run Transformations", [
    "dbt project: banking_dbt/",
    "Profile in banking_dbt/.dbt/profiles.yml → connects to Snowflake",
    "",
    "Layers:",
    "    Sources (sources.yml) → Points to BANKING.RAW tables",
    "    Staging (views) → stg_customers, stg_accounts, stg_transactions",
    "        • Extracts fields from VARIANT: v:id::string, v:email::string",
    "        • Deduplicates using ROW_NUMBER() OVER (PARTITION BY id ORDER BY created_at DESC)",
    "    Marts (tables) → dim_customers, dim_accounts, fact_transactions",
    "        • Dimensions built from SCD2 snapshots (effective_from, effective_to, is_current)",
    "        • Fact table uses incremental materialization with unique_key",
], "Commands: dbt debug  →  dbt run")

# ============================================================
# SLIDE 12: Step 12 — SCD2 Snapshots
# ============================================================
section_slide("12", "Run SCD2 Snapshots DAG", [
    "DAG: SCD2_snapshots (Schedule: @daily)",
    "File: docker/dags/scd_snapshots.py",
    "",
    "Task 1 — dbt_snapshot:",
    "    → Runs: dbt snapshot --profiles-dir /home/airflow/.dbt",
    "    → Compares current staging data vs. previous snapshot",
    "    → Detects changes and creates historical records",
    "",
    "Task 2 — dbt_run_marts:",
    "    → Runs: dbt run --select marts",
    "    → Rebuilds dimension and fact tables with latest snapshot data",
    "",
    "SCD Type-2 strategy: 'check' — monitors specific columns for changes:",
    "    • customers_snapshot checks: first_name, last_name, email",
    "    • accounts_snapshot checks: customer_id, account_type, balance",
], "Changed rows get dbt_valid_to timestamp; new row inserted with dbt_valid_from = now()")

# ============================================================
# SLIDE 13: How SCD2 Works
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text_box(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), "How SCD Type-2 Works", 36, ACCENT_BLUE, True)
items = [
    "Slowly Changing Dimension Type 2 = Full History Tracking",
    "",
    "Example: Customer 1 changes their email address",
    "",
    "  customer_id │ email              │ effective_from │ effective_to │ is_current",
    "  ────────────┼────────────────────┼────────────────┼──────────────┼───────────",
    "       1      │ old@email.com      │ 2026-01-01     │ 2026-05-15   │   FALSE",
    "       1      │ new@email.com      │ 2026-05-15     │     NULL     │   TRUE",
    "",
    "• Old row: dbt_valid_to gets set to the change timestamp",
    "• New row: inserted with dbt_valid_from = now, dbt_valid_to = NULL",
    "• is_current = TRUE only for the latest version",
    "",
    "This preserves complete audit history — you can query any point in time",
]
add_bullet_slide(s, Inches(0.5), Inches(1.4), Inches(12), Inches(5.5), items, 16, LIGHT_GRAY)

# ============================================================
# SLIDE 14: Step 13 — Verify in Snowflake
# ============================================================
section_slide("13", "Verify Data in Snowflake", [
    "Check raw data landed:",
    "    SELECT COUNT(*) FROM banking.raw.customers;",
    "    SELECT COUNT(*) FROM banking.raw.accounts;",
    "    SELECT COUNT(*) FROM banking.raw.transactions;",
    "",
    "Check staging views:",
    "    SELECT * FROM banking.analytics.stg_customers LIMIT 10;",
    "",
    "Check SCD2 snapshot history:",
    "    SELECT * FROM banking.analytics.customers_snapshot",
    "    WHERE customer_id = '1' ORDER BY dbt_valid_from;",
    "",
    "Check marts (dimension & fact tables):",
    "    SELECT * FROM banking.analytics.dim_customers WHERE is_current = TRUE LIMIT 10;",
    "    SELECT * FROM banking.analytics.fact_transactions LIMIT 10;",
])

# ============================================================
# SLIDE 15: Docker Services Summary
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text_box(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), "Docker Services Summary", 36, ACCENT_BLUE, True)
items = [
    "Service              │ Image                      │ Port        │ Purpose",
    "─────────────────────┼────────────────────────────┼─────────────┼──────────────────────",
    "postgres             │ postgres:15                │ 5432        │ Banking OLTP database",
    "minio                │ minio/minio:latest         │ 9000/9001   │ S3-compatible object storage",
    "zookeeper            │ confluentinc/cp-zookeeper  │ 2181        │ Kafka dependency",
    "kafka                │ confluentinc/cp-kafka      │ 9092/29092  │ Event streaming",
    "connect              │ debezium/connect:2.2       │ 8083        │ CDC connector",
    "airflow-webserver    │ Custom (apache/airflow)    │ 8080        │ Airflow UI",
    "airflow-scheduler    │ Custom (apache/airflow)    │ —           │ DAG executor",
    "airflow-postgres     │ postgres:15                │ 5433        │ Airflow metadata DB",
]
add_bullet_slide(s, Inches(0.3), Inches(1.5), Inches(12.5), Inches(5), items, 15, LIGHT_GRAY)
add_text_box(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5), "Start all: docker-compose up -d  |  Network: banking-mds-net", 14, MED_GRAY)

# ============================================================
# SLIDE 16: Key Files Reference
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text_box(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), "Key Files Reference", 36, ACCENT_BLUE, True)
items = [
    "postgres/schema.sql                           — Database table definitions",
    "data-generator/faker_generator.py             — Synthetic data generator",
    "consumer/postgres_to_minio.py                 — Direct PostgreSQL → MinIO transfer",
    "consumer/kafka_to_minio.py                    — Kafka CDC → MinIO consumer",
    "kafka-debezium/generate_and_post_connector.py — Debezium connector setup",
    "docker/dags/minio_to_snowflake_dag.py         — Airflow DAG: MinIO → Snowflake",
    "docker/dags/scd_snapshots.py                  — Airflow DAG: dbt snapshots + marts",
    "snowflake/queries.txt                         — Initial Snowflake setup queries",
    "banking_dbt/models/staging/stg_*.sql          — dbt staging models",
    "banking_dbt/models/marts/                     — dbt dimension & fact tables",
    "banking_dbt/snapshots/                        — dbt SCD2 snapshot definitions",
    "banking_dbt/.dbt/profiles.yml                 — dbt Snowflake connection profile",
    "docker-compose.yml                            — All Docker service definitions",
    "dockerfile-airflow.dockerfile                 — Custom Airflow image with dbt",
]
add_bullet_slide(s, Inches(0.3), Inches(1.4), Inches(12.5), Inches(5.5), items, 15, LIGHT_GRAY)

# ============================================================
# SLIDE 17: Thank You
# ============================================================
s = title_slide("Thank You!", "Questions?")

# ============================================================
# SAVE
# ============================================================
out_path = os.path.join(os.path.dirname(__file__), "..", "presentation", "Banking_Modern_Data_Stack_Training.pptx")
prs.save(out_path)
print(f"Saved to {out_path}")
